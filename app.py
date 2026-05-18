from __future__ import annotations

import os
import re
import shutil
import subprocess
import tempfile
import textwrap
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import streamlit as st
from openai import OpenAI, OpenAIError
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import imageio_ffmpeg


PAGE_TITLE = "PPT to Video Generator"
VIDEO_SIZE = (1280, 720)
MAX_SLIDES = 12
MAX_TTS_CHARS = 1800
DEFAULT_SECONDS_PER_SLIDE = 5.0
MIN_SECONDS_PER_SLIDE = 3.0
MAX_SECONDS_PER_SLIDE = 15.0


@dataclass
class SlideAsset:
    number: int
    title: str
    bullets: list[str]
    raw_text: str
    narration: str
    image_path: Path
    audio_path: Path | None = None
    duration: float = DEFAULT_SECONDS_PER_SLIDE


def get_openai_api_key() -> str | None:
    """Read the API key from Streamlit secrets first, then the environment."""
    return st.secrets.get("OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY")


def get_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def normalize_whitespace(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return "\n".join(normalize_whitespace(line) for line in shape.text.splitlines() if normalize_whitespace(line))


def extract_text_from_group(shape) -> list[str]:
    values: list[str] = []
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            values.extend(extract_text_from_group(child))
    else:
        text = shape_text(shape)
        if text:
            values.append(text)
    return values


def extract_slide_text(uploaded_pptx: Path, output_dir: Path, max_slides: int) -> list[SlideAsset]:
    prs = Presentation(uploaded_pptx)
    slides: list[SlideAsset] = []

    for index, slide in enumerate(prs.slides, start=1):
        if index > max_slides:
            break
        text_blocks: list[str] = []
        for shape in slide.shapes:
            text_blocks.extend(extract_text_from_group(shape))

        cleaned_blocks = [block for block in text_blocks if block]
        title = cleaned_blocks[0] if cleaned_blocks else f"Slide {index}"
        bullets = []
        for block in cleaned_blocks[1:]:
            bullets.extend(line for line in block.splitlines() if line.strip())
        if not bullets and len(cleaned_blocks) == 1:
            bullets = [cleaned_blocks[0]]

        raw_text = "\n".join(cleaned_blocks).strip() or f"Slide {index}"
        image_path = output_dir / f"slide_{index:02d}.png"
        slides.append(
            SlideAsset(
                number=index,
                title=title[:120],
                bullets=bullets[:8],
                raw_text=raw_text,
                narration=raw_text,
                image_path=image_path,
            )
        )

    if not slides:
        raise ValueError("The presentation does not contain any slides.")

    return slides


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    wrapped: list[str] = []
    for paragraph in text.splitlines() or [text]:
        words = paragraph.split()
        if not words:
            wrapped.append("")
            continue
        line = words[0]
        for word in words[1:]:
            candidate = f"{line} {word}"
            if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
                line = candidate
            else:
                wrapped.append(line)
                line = word
        wrapped.append(line)
    return wrapped


def draw_slide_image(slide: SlideAsset, slide_count: int) -> None:
    width, height = VIDEO_SIZE
    image = Image.new("RGB", VIDEO_SIZE, color=(245, 247, 251))
    draw = ImageDraw.Draw(image)

    title_font = get_font(48, bold=True)
    body_font = get_font(30)
    footer_font = get_font(20)

    # Lightweight branded layout with no ImageMagick/LibreOffice dependency.
    draw.rectangle((0, 0, width, 92), fill=(30, 64, 175))
    draw.rectangle((0, height - 58, width, height), fill=(226, 232, 240))
    draw.text((56, 24), f"Slide {slide.number}", font=footer_font, fill=(219, 234, 254))

    title_lines = wrap_text(draw, slide.title or f"Slide {slide.number}", title_font, width - 112)[:2]
    y = 128
    for line in title_lines:
        draw.text((56, y), line, font=title_font, fill=(15, 23, 42))
        y += 58

    y += 24
    bullet_source = slide.bullets or [slide.raw_text]
    for bullet in bullet_source[:7]:
        bullet = normalize_whitespace(bullet)
        if not bullet:
            continue
        lines = wrap_text(draw, bullet, body_font, width - 170)
        if y + (len(lines) * 38) > height - 100:
            draw.text((88, y), "…", font=body_font, fill=(51, 65, 85))
            break
        draw.ellipse((58, y + 12, 70, y + 24), fill=(37, 99, 235))
        for offset, line in enumerate(lines[:3]):
            draw.text((88, y + offset * 38), line, font=body_font, fill=(51, 65, 85))
        y += max(46, len(lines[:3]) * 38 + 16)

    footer = f"{slide.number} / {slide_count}"
    footer_width = draw.textbbox((0, 0), footer, font=footer_font)[2]
    draw.text((width - footer_width - 56, height - 40), footer, font=footer_font, fill=(71, 85, 105))
    image.save(slide.image_path, optimize=True)


def build_narration(client: OpenAI, slides: list[SlideAsset], model: str) -> None:
    for slide in slides:
        prompt = textwrap.dedent(
            f"""
            Write concise voiceover narration for one presentation slide.
            Requirements:
            - 2 to 4 short sentences
            - Natural spoken language
            - Do not mention bullet numbers or visual formatting
            - Keep under {MAX_TTS_CHARS} characters

            Slide text:
            {slide.raw_text[:3000]}
            """
        ).strip()
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "You convert slide text into clear, concise narration."},
                {"role": "user", "content": prompt},
            ],
            temperature=0.4,
            max_tokens=220,
        )
        narration = response.choices[0].message.content or slide.raw_text
        slide.narration = narration.strip()[:MAX_TTS_CHARS]


def estimate_duration_from_text(text: str) -> float:
    word_count = max(1, len(text.split()))
    return min(MAX_SECONDS_PER_SLIDE, max(MIN_SECONDS_PER_SLIDE, word_count / 2.4))



def audio_duration_seconds(audio_path: Path) -> float | None:
    try:
        completed = subprocess.run(
            [ffmpeg_executable(), "-i", str(audio_path), "-f", "null", "-"],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            check=False,
        )
    except Exception:
        return None

    match = re.search(r"Duration: (\d+):(\d+):(\d+(?:\.\d+)?)", completed.stderr)
    if not match:
        return None
    hours, minutes, seconds = match.groups()
    duration = int(hours) * 3600 + int(minutes) * 60 + float(seconds)
    return max(MIN_SECONDS_PER_SLIDE, min(MAX_SECONDS_PER_SLIDE, duration))

def generate_openai_tts(client: OpenAI, slides: list[SlideAsset], voice: str, tts_model: str, output_dir: Path) -> None:
    for slide in slides:
        audio_path = output_dir / f"audio_{slide.number:02d}.mp3"
        with client.audio.speech.with_streaming_response.create(
            model=tts_model,
            voice=voice,
            input=slide.narration[:MAX_TTS_CHARS],
            instructions="Speak clearly at a calm presenter pace.",
            response_format="mp3",
        ) as response:
            response.stream_to_file(audio_path)

        slide.audio_path = audio_path
        duration = audio_duration_seconds(audio_path)
        slide.duration = duration if duration else estimate_duration_from_text(slide.narration)


def ffmpeg_executable() -> str:
    try:
        return imageio_ffmpeg.get_ffmpeg_exe()
    except Exception:
        found = shutil.which("ffmpeg")
        if found:
            return found
        raise RuntimeError("FFmpeg is unavailable. Install imageio-ffmpeg or add FFmpeg to the runtime.")


def write_concat_file(path: Path, entries: Iterable[Path], durations: Iterable[float] | None = None) -> None:
    durations_list = list(durations) if durations is not None else None
    entries_list = list(entries)
    with path.open("w", encoding="utf-8") as file:
        for idx, entry in enumerate(entries_list):
            file.write(f"file '{entry.as_posix()}'\n")
            if durations_list is not None:
                file.write(f"duration {durations_list[idx]:.3f}\n")
        if durations_list is not None and entries_list:
            file.write(f"file '{entries_list[-1].as_posix()}'\n")


def run_ffmpeg(args: list[str]) -> None:
    completed = subprocess.run(args, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=False)
    if completed.returncode != 0:
        raise RuntimeError(completed.stderr[-3000:] or "FFmpeg failed without an error message.")


def create_video(slides: list[SlideAsset], output_path: Path, working_dir: Path) -> Path:
    ffmpeg = ffmpeg_executable()
    image_concat = working_dir / "images.txt"
    audio_concat = working_dir / "audio.txt"
    audio_path = working_dir / "narration.mp3"
    silent_audio = working_dir / "silence.mp3"
    video_only = working_dir / "video_only.mp4"

    write_concat_file(image_concat, [slide.image_path for slide in slides], [slide.duration for slide in slides])

    audio_files = [slide.audio_path for slide in slides if slide.audio_path]
    if audio_files:
        write_concat_file(audio_concat, audio_files)
        run_ffmpeg([ffmpeg, "-y", "-f", "concat", "-safe", "0", "-i", str(audio_concat), "-c", "copy", str(audio_path)])
    else:
        total_duration = sum(slide.duration for slide in slides)
        run_ffmpeg([
            ffmpeg,
            "-y",
            "-f",
            "lavfi",
            "-i",
            "anullsrc=channel_layout=stereo:sample_rate=44100",
            "-t",
            f"{total_duration:.3f}",
            "-q:a",
            "9",
            "-acodec",
            "libmp3lame",
            str(silent_audio),
        ])
        audio_path = silent_audio

    run_ffmpeg([
        ffmpeg,
        "-y",
        "-f",
        "concat",
        "-safe",
        "0",
        "-i",
        str(image_concat),
        "-fps_mode",
        "vfr",
        "-pix_fmt",
        "yuv420p",
        "-vf",
        "scale=1280:720:force_original_aspect_ratio=decrease,pad=1280:720:(ow-iw)/2:(oh-ih)/2",
        "-c:v",
        "libx264",
        "-preset",
        "veryfast",
        "-crf",
        "28",
        str(video_only),
    ])

    run_ffmpeg([
        ffmpeg,
        "-y",
        "-i",
        str(video_only),
        "-i",
        str(audio_path),
        "-c:v",
        "copy",
        "-c:a",
        "aac",
        "-b:a",
        "96k",
        "-shortest",
        "-movflags",
        "+faststart",
        str(output_path),
    ])
    return output_path


def save_upload(uploaded_file, output_dir: Path) -> Path:
    safe_name = re.sub(r"[^A-Za-z0-9_.-]", "_", uploaded_file.name or "presentation.pptx")
    ppt_path = output_dir / f"{uuid.uuid4().hex}_{safe_name}"
    ppt_path.write_bytes(uploaded_file.getvalue())
    return ppt_path


def render_app() -> None:
    st.set_page_config(page_title=PAGE_TITLE, page_icon="🎥", layout="wide")
    st.title("🎥 AI PPT-to-Video Generator")
    st.caption("Streamlit Cloud-only app: upload a .pptx, generate OpenAI narration, render a lightweight MP4, and download it.")

    with st.sidebar:
        st.header("Settings")
        max_slides = st.slider("Maximum slides", min_value=1, max_value=MAX_SLIDES, value=6)
        narration_model = st.selectbox("Narration model", options=["gpt-4o", "gpt-4o-mini"], index=0)
        tts_model = st.selectbox("TTS model", options=["gpt-4o-mini-tts", "tts-1"], index=0)
        voice = st.selectbox("Voice", options=["alloy", "coral", "marin", "cedar", "nova", "shimmer", "echo", "fable", "onyx", "sage"], index=0)
        st.info("Add OPENAI_API_KEY in Streamlit Community Cloud secrets before deploying.")

    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"], accept_multiple_files=False)
    generate = st.button("Generate video", type="primary", disabled=uploaded_file is None)

    if not generate:
        st.write("Upload a PowerPoint file to begin. The app renders simple readable slide images instead of relying on LibreOffice or ImageMagick.")
        return

    if uploaded_file is None:
        st.warning("Please upload a .pptx file first.")
        return

    api_key = get_openai_api_key()
    if not api_key:
        st.error("Missing OPENAI_API_KEY. Set it in Streamlit secrets or as an environment variable.")
        return

    progress = st.progress(0, text="Starting...")
    status = st.empty()

    try:
        with tempfile.TemporaryDirectory(prefix="ppt_video_") as tmp:
            workdir = Path(tmp)
            client = OpenAI(api_key=api_key)

            progress.progress(8, text="Saving upload...")
            ppt_path = save_upload(uploaded_file, workdir)

            progress.progress(20, text="Extracting slide text...")
            slides = extract_slide_text(ppt_path, workdir, max_slides=max_slides)
            status.write(f"Found {len(slides)} slide(s).")

            progress.progress(35, text="Creating readable slide images...")
            for slide in slides:
                draw_slide_image(slide, len(slides))

            progress.progress(50, text="Generating narration with OpenAI...")
            build_narration(client, slides, narration_model)

            progress.progress(68, text="Generating OpenAI TTS audio...")
            generate_openai_tts(client, slides, voice, tts_model, workdir)

            progress.progress(84, text="Rendering lightweight MP4...")
            output_path = workdir / "ppt_video.mp4"
            create_video(slides, output_path, workdir)

            video_bytes = output_path.read_bytes()
            progress.progress(100, text="Video ready.")
            st.success("Video generated successfully.")
            st.video(video_bytes)
            st.download_button(
                "Download MP4",
                data=video_bytes,
                file_name="ppt-video.mp4",
                mime="video/mp4",
                type="primary",
            )

            with st.expander("Generated narration"):
                for slide in slides:
                    st.markdown(f"**Slide {slide.number}: {slide.title}**")
                    st.write(slide.narration)
    except OpenAIError as exc:
        st.error(f"OpenAI request failed: {exc}")
    except subprocess.SubprocessError as exc:
        st.error(f"Video rendering failed: {exc}")
    except Exception as exc:
        st.error(f"Unable to generate video: {exc}")


if __name__ == "__main__":
    render_app()
